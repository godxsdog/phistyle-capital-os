#!/usr/bin/env python3
from __future__ import annotations

from datetime import datetime
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
import base64
import json
import re
from pathlib import Path
import ssl
import tempfile
import threading
import urllib.request
from urllib.parse import parse_qs
from urllib.parse import quote
from urllib.parse import urlencode
from urllib.parse import urlparse
from uuid import uuid4

import numpy as np
from PIL import Image, ImageEnhance, ImageFilter, ImageOps, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    import cv2
except Exception:  # pragma: no cover - optional local enhancement
    cv2 = None


APP_DIR = Path(__file__).resolve().parent
PROJECT_DIR = APP_DIR.parents[1]
DATA_DIR = PROJECT_DIR / "data" / "dental_cases"
CERT_FILE = PROJECT_DIR / "certs" / "points-wallet.crt"
KEY_FILE = PROJECT_DIR / "certs" / "points-wallet.key"

HTTPS_PORT = 8791
HTTP_PORT = 8790

CATEGORIES = {
    "intraoral": "口內照片",
    "extraoral": "口外照片",
    "xray": "X 光片",
}


class DentalHandler(SimpleHTTPRequestHandler):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=str(APP_DIR), **kwargs)

    def do_GET(self) -> None:
        path = urlparse(self.path).path
        if path == "/api/health":
            self.write_json({"ok": True, "app": "dental-ppt", "https": True})
            return
        if path == "/api/papers":
            params = self.query_params()
            query = params.get("q", [""])[0]
            limit = int(params.get("limit", ["8"])[0] or "8")
            try:
                self.write_json({"query": query, "papers": search_pubmed(query, limit=limit)})
            except Exception as exc:
                self.write_json({"query": query, "papers": [], "error": str(exc)}, status=502)
            return
        super().do_GET()

    def do_POST(self) -> None:
        path = urlparse(self.path).path
        if path == "/api/presentations":
            try:
                payload = self.read_json()
                case_dir = save_case(payload)
                ppt_path = build_presentation(case_dir, payload)
                self.write_file(
                    ppt_path,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    f"{safe_name(payload.get('caseTitle') or 'dental-case')}.pptx",
                )
            except Exception as exc:
                self.write_json({"ok": False, "error": str(exc)}, status=500)
            return
        self.send_error(404)

    def read_json(self) -> dict:
        length = int(self.headers.get("Content-Length", "0"))
        raw = self.rfile.read(length).decode("utf-8")
        return json.loads(raw or "{}")

    def query_params(self) -> dict:
        return parse_qs(urlparse(self.path).query)

    def write_json(self, payload: dict, status: int = 200) -> None:
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def write_file(self, path: Path, content_type: str, filename: str) -> None:
        data = path.read_bytes()
        self.send_response(200)
        self.send_header("Content-Type", content_type)
        ascii_name = re.sub(r"[^A-Za-z0-9._-]+", "-", filename).strip("-.")
        if not ascii_name or ascii_name.lower() == "pptx":
            ascii_name = "dental-case.pptx"
        encoded_name = quote(filename)
        self.send_header("Content-Disposition", f"attachment; filename=\"{ascii_name}\"; filename*=UTF-8''{encoded_name}")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)


class RedirectHandler(SimpleHTTPRequestHandler):
    def do_GET(self) -> None:
        self.redirect_to_https()

    def do_HEAD(self) -> None:
        self.redirect_to_https()

    def do_POST(self) -> None:
        self.redirect_to_https()

    def redirect_to_https(self) -> None:
        host = self.headers.get("Host", "127.0.0.1").split(":")[0]
        self.send_response(308)
        self.send_header("Location", f"https://{host}:{HTTPS_PORT}{self.path}")
        self.send_header("Cache-Control", "no-store")
        self.end_headers()


def safe_name(value: str) -> str:
    cleaned = re.sub(r"[^\w\u4e00-\u9fff.-]+", "-", value.strip(), flags=re.UNICODE)
    return cleaned.strip("-")[:80] or "dental-case"


def search_pubmed(query: str, limit: int = 8) -> list[dict]:
    query = (query or "").strip()
    if not query:
        return []
    limit = max(1, min(limit, 20))
    base = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils"
    search_params = urlencode({
        "db": "pubmed",
        "term": query,
        "retmode": "json",
        "retmax": str(limit),
        "sort": "relevance",
        "tool": "KaiChangDentalPresenter",
        "email": "kaichang@example.local",
    })
    search_url = f"{base}/esearch.fcgi?{search_params}"
    with urllib.request.urlopen(search_url, timeout=12) as response:
        search_data = json.loads(response.read().decode("utf-8"))
    pmids = search_data.get("esearchresult", {}).get("idlist", [])
    if not pmids:
        return []

    summary_params = urlencode({
        "db": "pubmed",
        "id": ",".join(pmids),
        "retmode": "json",
        "tool": "KaiChangDentalPresenter",
        "email": "kaichang@example.local",
    })
    summary_url = f"{base}/esummary.fcgi?{summary_params}"
    with urllib.request.urlopen(summary_url, timeout=12) as response:
        summary_data = json.loads(response.read().decode("utf-8"))

    result = summary_data.get("result", {})
    papers = []
    for pmid in pmids:
        item = result.get(pmid, {})
        if not item:
            continue
        authors = [author.get("name", "") for author in item.get("authors", [])[:3] if author.get("name")]
        year = (item.get("pubdate") or "").split(" ")[0]
        doi = ""
        for article_id in item.get("articleids", []):
            if article_id.get("idtype") == "doi":
                doi = article_id.get("value", "")
                break
        papers.append({
            "pmid": pmid,
            "title": clean_pubmed_text(item.get("title") or ""),
            "journal": item.get("fulljournalname") or item.get("source") or "",
            "year": year,
            "authors": authors,
            "doi": doi,
            "url": f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
        })
    return papers


def clean_pubmed_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip().rstrip(".")


def save_case(payload: dict) -> Path:
    created = datetime.now().strftime("%Y%m%d-%H%M%S")
    case_slug = safe_name(payload.get("caseTitle") or payload.get("patientCode") or "case")
    case_dir = DATA_DIR / f"{created}-{case_slug}-{uuid4().hex[:6]}"
    image_dir = case_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)

    clean_payload = {key: value for key, value in payload.items() if key != "images"}
    saved_images = []
    for item in payload.get("images", []):
        saved = save_image_item(image_dir, item)
        if saved:
            saved_images.append(saved)
    clean_payload["images"] = saved_images
    (case_dir / "case.json").write_text(json.dumps(clean_payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return case_dir


def save_image_item(image_dir: Path, item: dict) -> dict | None:
    data_url = item.get("dataUrl", "")
    if "," not in data_url:
        return None
    header, encoded = data_url.split(",", 1)
    ext = "png"
    if "jpeg" in header or "jpg" in header:
        ext = "jpg"
    elif "webp" in header:
        ext = "webp"
    name = safe_name(item.get("name") or item.get("slot") or "image")
    filename = f"{item.get('stage', 'case')}-{item.get('category', 'image')}-{name}-{uuid4().hex[:6]}.{ext}"
    raw_path = image_dir / filename
    raw_path.write_bytes(base64.b64decode(encoded))
    normalized_path = normalize_image(raw_path)
    return {
        "stage": item.get("stage", "before"),
        "category": item.get("category", "intraoral"),
        "slot": item.get("slot", ""),
        "name": item.get("name", filename),
        "path": str(normalized_path.relative_to(image_dir.parent)),
    }


def normalize_image(path: Path) -> Path:
    normalized = path.with_suffix(".jpg")
    with Image.open(path) as image:
        image = ImageOps.exif_transpose(image)
        if image.mode not in ("RGB", "L"):
            image = image.convert("RGB")
        image.thumbnail((2400, 1600), Image.Resampling.LANCZOS)
        image.save(normalized, "JPEG", quality=92, optimize=True)
    if normalized != path and path.exists():
        path.unlink()
    return normalized


def build_presentation(case_dir: Path, payload: dict) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    apply_theme(prs)

    title = payload.get("caseTitle") or "Dental Case Presentation"
    subtitle = f"{payload.get('patientCode') or 'Patient'} · {payload.get('treatment') or 'Treatment Plan'}"
    add_title_slide(prs, title, subtitle, payload)
    add_info_slide(prs, payload)

    image_items = load_saved_images(case_dir)
    missing_categories = []
    for category, label in CATEGORIES.items():
        pairs = pair_images([item for item in image_items if item["category"] == category])
        if pairs:
            add_category_overview(prs, label, pairs)
            for index, pair in enumerate(pairs, 1):
                add_comparison_slide(prs, payload, label, index, pair)
        else:
            missing_categories.append(label)

    if missing_categories:
        add_missing_categories_slide(prs, missing_categories)
    add_notes_slide(prs, payload)
    if payload.get("papers"):
        add_references_slide(prs, payload.get("papers", []))
    output = case_dir / f"{safe_name(title)}.pptx"
    prs.save(output)
    return output


def apply_theme(prs: Presentation) -> None:
    for layout in prs.slide_layouts:
        layout.name


def add_title_slide(prs: Presentation, title: str, subtitle: str, payload: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(12, 50, 54)
    add_text(slide, title, 0.75, 1.2, 11.8, 0.7, 34, RGBColor(255, 255, 255), bold=True)
    add_text(slide, subtitle, 0.78, 2.05, 11.0, 0.4, 17, RGBColor(196, 229, 224))
    meta = [
        ("醫師", payload.get("doctor") or "-"),
        ("日期", payload.get("caseDate") or datetime.now().strftime("%Y-%m-%d")),
        ("分類", payload.get("caseType") or "術前術後比對"),
    ]
    x = 0.78
    for label, value in meta:
        add_text(slide, label, x, 5.85, 1.3, 0.25, 10, RGBColor(164, 199, 194))
        add_text(slide, value, x, 6.18, 3.0, 0.35, 14, RGBColor(255, 255, 255), bold=True)
        x += 3.3


def add_info_slide(prs: Presentation, payload: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Case Summary")
    fields = [
        ("患者代碼", payload.get("patientCode")),
        ("治療項目", payload.get("treatment")),
        ("主訴 / 需求", payload.get("chiefConcern")),
        ("臨床備註", payload.get("clinicalNotes")),
    ]
    y = 1.25
    for label, value in fields:
        add_text(slide, label, 0.75, y, 2.0, 0.35, 12, RGBColor(49, 86, 88), bold=True)
        add_text(slide, value or "-", 2.55, y, 9.8, 0.42, 14, RGBColor(18, 30, 32))
        y += 0.72


def add_category_overview(prs: Presentation, label: str, pairs: list[dict]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, f"{label} · 術前術後對照")
    add_text(slide, f"共 {len(pairs)} 組可比對影像", 0.75, 1.1, 8, 0.35, 16, RGBColor(18, 30, 32), bold=True)
    y = 1.75
    for index, pair in enumerate(pairs[:8], 1):
        before_item = pair.get("before") or {}
        after_item = pair.get("after") or {}
        before = before_item.get("name") or "未提供"
        after = after_item.get("name") or "未提供"
        add_text(slide, f"{index}. {before}  →  {after}", 0.9, y, 11.4, 0.32, 12, RGBColor(38, 61, 62))
        y += 0.42


def add_comparison_slide(prs: Presentation, payload: dict, label: str, index: int, pair: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, f"{label} #{index}")
    if "口外" in label:
        bx, by, bw, bh = 1.85, 1.22, 3.65, 5.35
        ax, ay, aw, ah = 7.85, 1.22, 3.65, 5.35
        note_y = 6.88
    else:
        bx, by, bw, bh = 0.75, 1.55, 5.85, 3.90
        ax, ay, aw, ah = 6.75, 1.55, 5.85, 3.90
        note_y = 5.95

    add_text(slide, "Before", bx, 1.02, bw, 0.3, 15, RGBColor(129, 59, 47), bold=True)
    add_text(slide, "After", ax, 1.02, aw, 0.3, 15, RGBColor(26, 105, 91), bold=True)
    add_pair_images_or_placeholders(slide, pair, bx, by, bw, bh, ax, ay, aw, ah)
    note = payload.get("comparisonNotes") or "自動排版：請由醫師確認影像方向、尺度與臨床判讀。"
    add_text(slide, f"{note} 影像已做同組位置、尺度、亮度、對比與色彩標準化。", 0.8, note_y, 11.7, 0.35, 9, RGBColor(83, 97, 99))


def add_notes_slide(prs: Presentation, payload: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Treatment Notes")
    text = payload.get("finalNotes") or "請補入治療重點、患者溝通內容與術後追蹤計畫。"
    add_text(slide, text, 0.85, 1.3, 11.6, 4.8, 18, RGBColor(18, 30, 32))


def add_missing_categories_slide(prs: Presentation, missing_categories: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "未提供影像分類")
    add_text(slide, "以下分類未上傳影像，本簡報仍可正常產生。", 0.85, 1.2, 11.4, 0.42, 16, RGBColor(18, 30, 32), bold=True)
    y = 1.95
    for label in missing_categories:
        add_text(slide, f"- {label}", 1.05, y, 10.5, 0.35, 14, RGBColor(83, 97, 99))
        y += 0.48


def add_references_slide(prs: Presentation, papers: list[dict]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Evidence / References")
    add_text(slide, "PubMed 文獻搜尋結果，請由醫師確認是否適合作為本案例佐證。", 0.75, 0.85, 11.8, 0.35, 11, RGBColor(83, 97, 99))
    y = 1.28
    for index, paper in enumerate(papers[:8], 1):
        authors = ", ".join(paper.get("authors", [])[:3])
        if len(paper.get("authors", [])) > 3:
            authors += " et al."
        year = paper.get("year") or "n.d."
        journal = paper.get("journal") or "PubMed"
        title = paper.get("title") or "Untitled"
        pmid = paper.get("pmid") or ""
        doi = f" DOI: {paper.get('doi')}" if paper.get("doi") else ""
        citation = f"{index}. {authors} ({year}). {title}. {journal}. PMID: {pmid}.{doi}"
        add_text(slide, citation, 0.85, y, 11.6, 0.54, 9, RGBColor(18, 30, 32))
        y += 0.62


def load_saved_images(case_dir: Path) -> list[dict]:
    payload = json.loads((case_dir / "case.json").read_text(encoding="utf-8"))
    images = []
    for item in payload.get("images", []):
        copied = dict(item)
        copied["absolute_path"] = case_dir / copied["path"]
        images.append(copied)
    return images


def pair_images(items: list[dict]) -> list[dict]:
    before = [item for item in items if item.get("stage") == "before"]
    after = [item for item in items if item.get("stage") == "after"]
    count = max(len(before), len(after))
    pairs = []
    for index in range(count):
        pairs.append({
            "before": before[index] if index < len(before) else None,
            "after": after[index] if index < len(after) else None,
        })
    return pairs


def add_header(slide, title: str) -> None:
    add_band(slide, 0, 0, 13.333, 0.62, RGBColor(12, 50, 54))
    add_text(slide, title, 0.55, 0.16, 11.5, 0.28, 15, RGBColor(255, 255, 255), bold=True)


def add_band(slide, x: float, y: float, w: float, h: float, color: RGBColor):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, color: RGBColor, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def add_image_or_placeholder(slide, item: dict | None, x: float, y: float, w: float, h: float) -> None:
    if not item:
        add_band(slide, x, y, w, h, RGBColor(238, 243, 242))
        add_text(slide, "未提供影像", x + 2.1, y + 2.05, 2.0, 0.35, 14, RGBColor(83, 97, 99), bold=True)
        return
    path = Path(item["absolute_path"])
    with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
        fitted = fit_image(path, Inches(w), Inches(h))
        fitted.save(tmp.name, "JPEG", quality=94)
        slide.shapes.add_picture(tmp.name, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    add_text(slide, item.get("name") or "", x, y + h + 0.08, w, 0.2, 8, RGBColor(83, 97, 99))


def add_pair_images_or_placeholders(slide, pair: dict, bx: float, by: float, bw: float, bh: float, ax: float, ay: float, aw: float, ah: float) -> None:
    before_item = pair.get("before")
    after_item = pair.get("after")
    before_image, after_image = prepare_pair_images(before_item, after_item, bw, bh)

    if before_image:
        add_prepared_image(slide, before_image, bx, by, bw, bh)
        add_text(slide, before_item.get("name") or "", bx, by + bh + 0.08, bw, 0.2, 8, RGBColor(83, 97, 99))
    else:
        add_image_or_placeholder(slide, None, bx, by, bw, bh)

    if after_image:
        add_prepared_image(slide, after_image, ax, ay, aw, ah)
        add_text(slide, after_item.get("name") or "", ax, ay + ah + 0.08, aw, 0.2, 8, RGBColor(83, 97, 99))
    else:
        add_image_or_placeholder(slide, None, ax, ay, aw, ah)


def add_prepared_image(slide, image: Image.Image, x: float, y: float, w: float, h: float) -> None:
    with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
        image.save(tmp.name, "JPEG", quality=94)
        slide.shapes.add_picture(tmp.name, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def prepare_pair_images(before_item: dict | None, after_item: dict | None, w: float, h: float) -> tuple[Image.Image | None, Image.Image | None]:
    target_ratio = float(w) / float(h)
    target_size = target_pixels(target_ratio)
    category = (before_item or after_item or {}).get("category", "")
    before_raw = load_rgb(before_item) if before_item else None
    after_raw = load_rgb(after_item) if after_item else None
    before_image = align_subject(before_raw, target_ratio, target_size, category) if before_raw else None
    after_image = align_subject(after_raw, target_ratio, target_size, category) if after_raw else None
    return normalize_pair_color_and_tone(before_image, after_image, category)


def target_pixels(target_ratio: float) -> tuple[int, int]:
    width = 1800
    height = max(1, round(width / target_ratio))
    return width, height


def load_rgb(item: dict) -> Image.Image:
    return ImageOps.exif_transpose(Image.open(Path(item["absolute_path"]))).convert("RGB")


def align_subject(image: Image.Image, target_ratio: float, target_size: tuple[int, int], category: str) -> Image.Image:
    if category == "extraoral":
        crop = extraoral_eye_aligned_crop(image, target_ratio)
        if crop is None:
            bbox = extraoral_subject_bbox(image)
            crop = crop_box_from_bbox(image.size, bbox, target_ratio, padding=1.20, min_fill=0.54, y_bias=-0.08)
    elif category == "intraoral":
        bbox = intraoral_subject_bbox(image)
        crop = crop_box_from_bbox(image.size, bbox, target_ratio, padding=1.04, min_fill=0.92, y_bias=0.0)
    else:
        bbox = salient_bbox(image)
        crop = crop_box_from_bbox(image.size, bbox, target_ratio, padding=1.12, min_fill=0.80, y_bias=0.0)
    return image.crop(crop).resize(target_size, Image.Resampling.LANCZOS)


def extraoral_eye_aligned_crop(image: Image.Image, target_ratio: float) -> tuple[int, int, int, int] | None:
    landmarks = detect_extraoral_eyes(image)
    if not landmarks:
        return None
    left_eye, right_eye = landmarks
    eye_mid_x = (left_eye[0] + right_eye[0]) / 2
    eye_mid_y = (left_eye[1] + right_eye[1]) / 2
    eye_distance = abs(right_eye[0] - left_eye[0])
    if eye_distance < image.width * 0.08:
        return None

    # Standardize the portrait: both eyes land at the same x/y and the same
    # inter-eye scale, so the nose and facial midline follow the same frame.
    crop_w = eye_distance / 0.305
    crop_h = crop_w / target_ratio
    crop_w = min(image.width, max(crop_w, image.width * 0.48))
    crop_h = min(image.height, max(crop_h, image.height * 0.56))
    if crop_w / crop_h > target_ratio:
        crop_w = crop_h * target_ratio
    else:
        crop_h = crop_w / target_ratio

    x1 = eye_mid_x - crop_w * 0.50
    y1 = eye_mid_y - crop_h * 0.385
    x1 = min(max(0, x1), image.width - crop_w)
    y1 = min(max(0, y1), image.height - crop_h)
    return (int(round(x1)), int(round(y1)), int(round(x1 + crop_w)), int(round(y1 + crop_h)))


def detect_extraoral_eyes(image: Image.Image) -> tuple[tuple[float, float], tuple[float, float]] | None:
    cv2_result = detect_extraoral_eyes_cv2(image)
    if cv2_result:
        return cv2_result

    scale_w = 700
    scale = scale_w / image.width
    scale_h = max(1, round(image.height * scale))
    work_rgb = image.convert("RGB").resize((scale_w, scale_h), Image.Resampling.BILINEAR)
    face_box = extraoral_skin_bbox_scaled(np.asarray(work_rgb, dtype=np.float32))
    work = ImageOps.autocontrast(work_rgb.convert("L"))
    arr = np.asarray(work, dtype=np.float32)

    if face_box:
        fx1, fy1, fx2, fy2 = face_box
        face_w = fx2 - fx1
        face_h = fy2 - fy1
        x1 = int(max(scale_w * 0.16, fx1 + face_w * 0.06))
        x2 = int(min(scale_w * 0.84, fx2 - face_w * 0.06))
        y1 = int(max(scale_h * 0.28, fy1 + face_h * 0.16))
        y2 = int(min(scale_h * 0.58, fy1 + face_h * 0.48))
    else:
        y1 = int(scale_h * 0.31)
        y2 = int(scale_h * 0.55)
        x1 = int(scale_w * 0.22)
        x2 = int(scale_w * 0.78)
    if y2 <= y1 or x2 <= x1:
        return None

    roi = arr[y1:y2, x1:x2]
    threshold = min(112.0, float(np.percentile(roi, 18)))
    darkness = np.clip(threshold - roi, 0, None)
    if float(darkness.max()) < 8:
        return None

    half = darkness.shape[1] // 2
    left = strongest_eye_window(darkness[:, :half], x1, y1, scale)
    right = strongest_eye_window(darkness[:, half:], x1 + half, y1, scale)
    if not left or not right:
        return None
    if abs(left[1] - right[1]) > image.height * 0.06:
        shared_y = (left[1] + right[1]) / 2
        left = (left[0], shared_y)
        right = (right[0], shared_y)
    if right[0] <= left[0]:
        return None
    return left, right


def detect_extraoral_eyes_cv2(image: Image.Image) -> tuple[tuple[float, float], tuple[float, float]] | None:
    if cv2 is None:
        return None
    gray = np.asarray(image.convert("L"))
    gray = cv2.equalizeHist(gray)
    face_path = str(Path(cv2.data.haarcascades) / "haarcascade_frontalface_default.xml")
    eye_path = str(Path(cv2.data.haarcascades) / "haarcascade_eye_tree_eyeglasses.xml")
    face_cascade = cv2.CascadeClassifier(face_path)
    eye_cascade = cv2.CascadeClassifier(eye_path)
    if face_cascade.empty() or eye_cascade.empty():
        return None

    min_face = max(120, int(min(image.size) * 0.28))
    faces = face_cascade.detectMultiScale(gray, scaleFactor=1.04, minNeighbors=4, minSize=(min_face, min_face))
    if len(faces) == 0:
        return None
    image_center = image.width / 2
    face = max(
        faces,
        key=lambda rect: rect[2] * rect[3] - abs((rect[0] + rect[2] / 2) - image_center) * 20,
    )
    x, y, w, h = [int(v) for v in face]
    roi_y1 = y + int(h * 0.18)
    roi_y2 = y + int(h * 0.56)
    roi_x1 = x + int(w * 0.08)
    roi_x2 = x + int(w * 0.92)
    roi = gray[roi_y1:roi_y2, roi_x1:roi_x2]
    eyes = eye_cascade.detectMultiScale(
        roi,
        scaleFactor=1.04,
        minNeighbors=3,
        minSize=(max(18, int(w * 0.08)), max(12, int(h * 0.045))),
    )
    candidates = []
    for ex, ey, ew, eh in eyes:
        cx = roi_x1 + ex + ew / 2
        cy = roi_y1 + ey + eh / 2
        if not (x + w * 0.18 <= cx <= x + w * 0.82):
            continue
        if not (y + h * 0.25 <= cy <= y + h * 0.52):
            continue
        candidates.append((float(cx), float(cy), float(ew * eh)))

    if len(candidates) >= 2:
        best_pair = None
        best_score = -1.0
        for i, left in enumerate(candidates):
            for right in candidates[i + 1:]:
                if right[0] < left[0]:
                    left, right = right, left
                dx = right[0] - left[0]
                dy = abs(right[1] - left[1])
                if dx < w * 0.18 or dx > w * 0.58 or dy > h * 0.08:
                    continue
                score = dx - dy * 3 + left[2] * 0.0005 + right[2] * 0.0005
                if score > best_score:
                    best_score = score
                    shared_y = (left[1] + right[1]) / 2
                    best_pair = ((left[0], shared_y), (right[0], shared_y))
        if best_pair:
            return best_pair

    # If eyes are not found, use a conservative face-derived approximation.
    return (
        (x + w * 0.34, y + h * 0.40),
        (x + w * 0.66, y + h * 0.40),
    )


def extraoral_skin_bbox_scaled(rgb: np.ndarray) -> tuple[int, int, int, int] | None:
    height, width, _ = rgb.shape
    r = rgb[:, :, 0]
    g = rgb[:, :, 1]
    b = rgb[:, :, 2]
    luminance = r * 0.299 + g * 0.587 + b * 0.114
    saturation = rgb.max(axis=2) - rgb.min(axis=2)
    mask = (
        (luminance > 95)
        & (luminance < 242)
        & (saturation > 8)
        & (saturation < 92)
        & (r > b + 8)
        & (g > b + 2)
        & (r > g * 0.92)
    )
    region = np.zeros_like(mask, dtype=bool)
    region[int(height * 0.22):int(height * 0.78), int(width * 0.16):int(width * 0.84)] = True
    mask &= region
    ys, xs = np.where(mask)
    if len(xs) < 500 or len(ys) < 500:
        return None
    left = int(np.percentile(xs, 2))
    right = int(np.percentile(xs, 98))
    top = int(np.percentile(ys, 2))
    bottom = int(np.percentile(ys, 98))
    if right - left < width * 0.18 or bottom - top < height * 0.18:
        return None
    return left, top, right, bottom


def strongest_eye_window(score: np.ndarray, offset_x: int, offset_y: int, scale: float) -> tuple[float, float] | None:
    if score.size == 0:
        return None
    window_w = max(16, round(score.shape[1] * 0.18))
    window_h = max(10, round(score.shape[0] * 0.14))
    if score.shape[0] <= window_h or score.shape[1] <= window_w:
        return None

    integral = np.pad(score, ((1, 0), (1, 0)), mode="constant").cumsum(axis=0).cumsum(axis=1)
    sums = (
        integral[window_h:, window_w:]
        - integral[:-window_h, window_w:]
        - integral[window_h:, :-window_w]
        + integral[:-window_h, :-window_w]
    )
    y, x = np.unravel_index(int(np.argmax(sums)), sums.shape)
    confidence = float(sums[y, x]) / (window_w * window_h)
    if confidence < 2.0:
        return None
    center_x = (offset_x + x + window_w / 2) / scale
    center_y = (offset_y + y + window_h / 2) / scale
    return center_x, center_y


def extraoral_subject_bbox(image: Image.Image) -> tuple[int, int, int, int]:
    small_w = 420
    small_h = max(1, round(small_w * image.height / image.width))
    work = image.resize((small_w, small_h), Image.Resampling.BILINEAR).convert("RGB")
    arr = np.asarray(work, dtype=np.float32)
    maxc = arr.max(axis=2)
    minc = arr.min(axis=2)
    saturation = maxc - minc
    luminance = arr[:, :, 0] * 0.299 + arr[:, :, 1] * 0.587 + arr[:, :, 2] * 0.114
    mask = ((saturation > 16) | (luminance < 220)) & (luminance < 248)

    height, width = mask.shape
    center_band = np.zeros_like(mask, dtype=bool)
    center_band[:, int(width * 0.12):int(width * 0.88)] = True
    mask &= center_band
    ys, xs = np.where(mask)
    if len(xs) < 200 or len(ys) < 200:
        return salient_bbox(image)

    scale_x = image.width / width
    scale_y = image.height / height
    left = int(max(0, np.percentile(xs, 1.0) * scale_x))
    right = int(min(image.width, (np.percentile(xs, 99.0) + 1) * scale_x))
    top = int(max(0, np.percentile(ys, 0.5) * scale_y))
    bottom = int(min(image.height, (np.percentile(ys, 93.0) + 1) * scale_y))
    return expand_bbox((left, top, right, bottom), image.size, 0.05, 0.05)


def intraoral_subject_bbox(image: Image.Image) -> tuple[int, int, int, int]:
    arr = np.asarray(image.resize((520, max(1, round(520 * image.height / image.width))), Image.Resampling.BILINEAR).convert("RGB"), dtype=np.float32)
    maxc = arr.max(axis=2)
    minc = arr.min(axis=2)
    saturation = maxc - minc
    luminance = arr[:, :, 0] * 0.299 + arr[:, :, 1] * 0.587 + arr[:, :, 2] * 0.114
    mask = (saturation > 18) & (luminance > 45) & (luminance < 245)
    ys, xs = np.where(mask)
    if len(xs) < 300 or len(ys) < 300:
        return (0, 0, image.width, image.height)
    scale_x = image.width / arr.shape[1]
    scale_y = image.height / arr.shape[0]
    left = int(max(0, np.percentile(xs, 0.5) * scale_x))
    right = int(min(image.width, (np.percentile(xs, 99.5) + 1) * scale_x))
    top = int(max(0, np.percentile(ys, 0.5) * scale_y))
    bottom = int(min(image.height, (np.percentile(ys, 99.5) + 1) * scale_y))
    return expand_bbox((left, top, right, bottom), image.size, 0.02, 0.02)


def salient_bbox(image: Image.Image) -> tuple[int, int, int, int]:
    work = ImageOps.autocontrast(image.convert("L")).resize((480, max(1, round(480 * image.height / image.width))), Image.Resampling.BILINEAR)
    edges = work.filter(ImageFilter.FIND_EDGES)
    arr = np.asarray(edges, dtype=np.float32)
    threshold = max(18.0, float(arr.mean() + arr.std() * 0.65))
    mask = arr > threshold
    ys, xs = np.where(mask)
    if len(xs) < 60 or len(ys) < 60:
        return (0, 0, image.width, image.height)
    scale_x = image.width / work.width
    scale_y = image.height / work.height
    left = int(max(0, xs.min() * scale_x))
    right = int(min(image.width, (xs.max() + 1) * scale_x))
    top = int(max(0, ys.min() * scale_y))
    bottom = int(min(image.height, (ys.max() + 1) * scale_y))
    min_w = image.width * 0.28
    min_h = image.height * 0.28
    if right - left < min_w or bottom - top < min_h:
        return (0, 0, image.width, image.height)
    return (left, top, right, bottom)


def expand_bbox(bbox: tuple[int, int, int, int], size: tuple[int, int], pad_x: float, pad_y: float) -> tuple[int, int, int, int]:
    image_w, image_h = size
    left, top, right, bottom = bbox
    width = right - left
    height = bottom - top
    return (
        max(0, int(round(left - width * pad_x))),
        max(0, int(round(top - height * pad_y))),
        min(image_w, int(round(right + width * pad_x))),
        min(image_h, int(round(bottom + height * pad_y))),
    )


def crop_box_from_bbox(
    size: tuple[int, int],
    bbox: tuple[int, int, int, int],
    target_ratio: float,
    padding: float = 1.2,
    min_fill: float = 0.45,
    y_bias: float = 0.0,
) -> tuple[int, int, int, int]:
    image_w, image_h = size
    left, top, right, bottom = bbox
    center_x = (left + right) / 2
    center_y = (top + bottom) / 2 + (bottom - top) * y_bias
    subject_w = max(1, right - left)
    subject_h = max(1, bottom - top)
    crop_w = subject_w * padding
    crop_h = subject_h * padding
    if crop_w / crop_h > target_ratio:
        crop_h = crop_w / target_ratio
    else:
        crop_w = crop_h * target_ratio
    crop_w = min(image_w, max(crop_w, image_w * min_fill))
    crop_h = min(image_h, max(crop_h, image_h * min_fill))
    if crop_w / crop_h > target_ratio:
        crop_w = crop_h * target_ratio
    else:
        crop_h = crop_w / target_ratio
    x1 = center_x - crop_w / 2
    y1 = center_y - crop_h / 2
    x1 = min(max(0, x1), image_w - crop_w)
    y1 = min(max(0, y1), image_h - crop_h)
    return (int(round(x1)), int(round(y1)), int(round(x1 + crop_w)), int(round(y1 + crop_h)))


def normalize_pair_color_and_tone(before_image: Image.Image | None, after_image: Image.Image | None, category: str) -> tuple[Image.Image | None, Image.Image | None]:
    images = [image for image in (before_image, after_image) if image is not None]
    if len(images) < 2:
        return before_image, after_image

    if category == "extraoral" and cv2 is not None:
        return normalize_extraoral_regions(before_image, after_image)

    channel_stats = [rgb_stats(image, category) for image in images]
    target_mean = np.mean([item[0] for item in channel_stats], axis=0)
    target_std = np.maximum(8.0, np.mean([item[1] for item in channel_stats], axis=0))

    return (
        match_rgb_stats(before_image, target_mean, target_std, category) if before_image else None,
        match_rgb_stats(after_image, target_mean, target_std, category) if after_image else None,
    )


def normalize_extraoral_regions(before_image: Image.Image | None, after_image: Image.Image | None) -> tuple[Image.Image | None, Image.Image | None]:
    images = [image for image in (before_image, after_image) if image is not None]
    luma_profiles = [extraoral_luma_profile(image) for image in images]
    target_luma_mean = float(np.mean([item["mean"] for item in luma_profiles]))
    target_luma_std = float(max(6.0, np.min([item["std"] for item in luma_profiles]) * 0.94))
    luma_normalized = [
        Image.fromarray(match_global_luma(np.asarray(image, dtype=np.float32), target_luma_mean, target_luma_std, strength=1.0).astype(np.uint8), "RGB")
        for image in images
    ]
    calibration_profiles = [extraoral_global_calibration_profile(image) for image in luma_normalized]
    target_calibration = {
        "mean": np.mean([profile["mean"] for profile in calibration_profiles], axis=0),
        "std": np.maximum(2.0, np.min([profile["std"] for profile in calibration_profiles], axis=0) * 0.94),
    }

    return (
        match_extraoral_global(before_image, target_calibration, target_luma_mean, target_luma_std) if before_image else None,
        match_extraoral_global(after_image, target_calibration, target_luma_mean, target_luma_std) if after_image else None,
    )


def extraoral_global_calibration_profile(image: Image.Image) -> dict[str, np.ndarray]:
    rgb = np.asarray(image, dtype=np.float32)
    masks = extraoral_region_masks(rgb)
    weighted_sources = []
    for zone, weight in (("skin", 4), ("hair", 3), ("gray", 1)):
        sample = rgb[masks[zone]]
        if len(sample) < 300:
            continue
        limit = min(len(sample), 120000)
        if len(sample) > limit:
            step = max(1, len(sample) // limit)
            sample = sample[::step][:limit]
        weighted_sources.append((sample, weight))
    if not weighted_sources:
        flat = rgb.reshape(-1, 3)
        return {"mean": flat.mean(axis=0), "std": np.maximum(1.0, flat.std(axis=0))}

    weighted_samples = []
    for sample, weight in weighted_sources:
        weighted_samples.extend([sample] * weight)
    calibration = np.vstack(weighted_samples)
    return {"mean": calibration.mean(axis=0), "std": np.maximum(1.0, calibration.std(axis=0))}


def match_extraoral_global(
    image: Image.Image,
    target_calibration: dict[str, np.ndarray],
    target_luma_mean: float,
    target_luma_std: float,
) -> Image.Image:
    adjusted = match_global_luma(np.asarray(image, dtype=np.float32), target_luma_mean, target_luma_std, strength=1.0)
    current_profile = extraoral_global_calibration_profile(Image.fromarray(np.clip(adjusted, 0, 255).astype(np.uint8), "RGB"))
    current_mean = current_profile["mean"]
    current_std = current_profile["std"]
    wanted_mean = target_calibration["mean"]
    wanted_std = target_calibration["std"]
    transformed = (adjusted - current_mean) * (wanted_std / np.maximum(1.0, current_std)) + wanted_mean
    adjusted = adjusted * 0.10 + transformed * 0.90
    adjusted = match_global_luma(adjusted, target_luma_mean, target_luma_std, strength=1.0)
    adjusted = np.clip(adjusted, 0, 255).astype(np.uint8)
    result = Image.fromarray(adjusted, "RGB")
    result = ImageEnhance.Contrast(result).enhance(1.01)
    return ImageEnhance.Sharpness(result).enhance(1.02)


def extraoral_luma_profile(image: Image.Image) -> dict[str, float]:
    rgb = np.asarray(image, dtype=np.float32)
    luma = rgb[:, :, 0] * 0.299 + rgb[:, :, 1] * 0.587 + rgb[:, :, 2] * 0.114
    return {"mean": float(luma.mean()), "std": float(max(1.0, luma.std()))}


def extraoral_color_profile(image: Image.Image) -> dict[str, dict[str, np.ndarray | int]]:
    rgb = np.asarray(image, dtype=np.float32)
    masks = extraoral_region_masks(rgb)
    profile = {}
    for zone, mask in masks.items():
        sample = rgb[mask]
        if len(sample) < 300:
            sample = rgb.reshape(-1, 3)
        profile[zone] = {
            "mean": sample.mean(axis=0),
            "std": np.maximum(1.0, sample.std(axis=0)),
            "count": int(len(sample)),
        }
    return profile


def match_extraoral_regions(
    image: Image.Image,
    target: dict[str, dict[str, np.ndarray]],
    target_luma_mean: float,
    target_luma_std: float,
) -> Image.Image:
    original = np.asarray(image, dtype=np.float32)
    masks = extraoral_region_masks(original)
    adjusted = match_global_luma(original, target_luma_mean, target_luma_std, strength=0.92)

    if "gray" in target:
        gray_mask = masks["gray"]
        if np.count_nonzero(gray_mask) > 300:
            gray_mean = original[gray_mask].mean(axis=0)
            target_gray = target["gray"]["mean"]
            gain = np.clip(target_gray / np.maximum(1.0, gray_mean), 0.75, 1.25)
            adjusted = adjusted * (gain * 0.70 + 0.30)

    for _ in range(2):
        masks = extraoral_region_masks(adjusted)
        for zone, strength in (("skin", 0.96), ("hair", 0.92), ("gray", 0.78)):
            if zone not in target:
                continue
            mask = masks[zone]
            if np.count_nonzero(mask) < 300:
                continue
            zone_mean = adjusted[mask].mean(axis=0)
            zone_std = np.maximum(1.0, adjusted[mask].std(axis=0))
            wanted_mean = target[zone]["mean"]
            wanted_std = target[zone]["std"]
            zone_adjusted = (adjusted - zone_mean) * (wanted_std / zone_std) + wanted_mean
            zone_adjusted = adjusted * (1.0 - strength) + zone_adjusted * strength
            alpha = soft_mask(mask)
            adjusted = adjusted * (1.0 - alpha) + zone_adjusted * alpha

    adjusted = match_global_luma(adjusted, target_luma_mean, target_luma_std, strength=1.0)
    adjusted = np.clip(adjusted, 0, 255).astype(np.uint8)
    result = Image.fromarray(adjusted, "RGB")
    result = ImageEnhance.Contrast(result).enhance(1.01)
    return ImageEnhance.Sharpness(result).enhance(1.02)


def match_global_luma(rgb: np.ndarray, target_mean: float, target_std: float, strength: float) -> np.ndarray:
    luma = rgb[:, :, 0] * 0.299 + rgb[:, :, 1] * 0.587 + rgb[:, :, 2] * 0.114
    current_mean = float(luma.mean())
    current_std = float(max(1.0, luma.std()))
    desired_luma = (luma - current_mean) * (target_std / current_std) + target_mean
    delta = (desired_luma - luma)[:, :, None]
    return np.clip(rgb + delta * strength, 0, 255)


def extraoral_region_masks(rgb: np.ndarray) -> dict[str, np.ndarray]:
    height, width, _ = rgb.shape
    r = rgb[:, :, 0]
    g = rgb[:, :, 1]
    b = rgb[:, :, 2]
    luminance = r * 0.299 + g * 0.587 + b * 0.114
    saturation = rgb.max(axis=2) - rgb.min(axis=2)
    center = np.zeros((height, width), dtype=bool)
    center[int(height * 0.10):int(height * 0.82), int(width * 0.13):int(width * 0.87)] = True
    face_core = np.zeros((height, width), dtype=bool)
    face_core[int(height * 0.22):int(height * 0.76), int(width * 0.18):int(width * 0.82)] = True
    upper = np.zeros((height, width), dtype=bool)
    upper[int(height * 0.02):int(height * 0.48), int(width * 0.08):int(width * 0.92)] = True
    border = np.ones((height, width), dtype=bool)
    border[int(height * 0.10):int(height * 0.86), int(width * 0.10):int(width * 0.90)] = False

    skin = (
        face_core
        & (luminance > 72)
        & (luminance < 246)
        & (saturation > 4)
        & (saturation < 112)
        & (r > b + 4)
        & (g > b - 6)
        & (r > g * 0.84)
    )
    hair = (
        upper
        & center
        & (luminance > 18)
        & (luminance < 145)
        & (saturation > 12)
        & (r > b - 4)
        & ~skin
    )
    gray = (
        ((border | ~center) | (luminance > 205))
        & (luminance > 118)
        & (luminance < 250)
        & (saturation < 32)
        & ~skin
        & ~hair
    )
    return {"skin": skin, "hair": hair, "gray": gray}


def soft_mask(mask: np.ndarray) -> np.ndarray:
    alpha = mask.astype(np.float32)
    if cv2 is not None:
        alpha = cv2.GaussianBlur(alpha, (0, 0), sigmaX=9, sigmaY=9)
    else:
        alpha = np.asarray(Image.fromarray((alpha * 255).astype(np.uint8)).filter(ImageFilter.GaussianBlur(9)), dtype=np.float32) / 255.0
    alpha = np.clip(alpha, 0.0, 1.0)
    return alpha[:, :, None]


def extraoral_lab_stats(image: Image.Image) -> tuple[np.ndarray, np.ndarray]:
    rgb = np.asarray(image, dtype=np.uint8)
    lab = cv2.cvtColor(rgb, cv2.COLOR_RGB2LAB).astype(np.float32)
    mask = extraoral_skin_mask(rgb.astype(np.float32))
    if np.count_nonzero(mask) < 400:
        height, width = mask.shape
        mask = np.zeros_like(mask, dtype=bool)
        mask[int(height * 0.20):int(height * 0.72), int(width * 0.18):int(width * 0.82)] = True
    sample = lab[mask]
    return sample.mean(axis=0), np.maximum(1.0, sample.std(axis=0))


def match_lab_stats(image: Image.Image, target_mean: np.ndarray, target_std: np.ndarray) -> Image.Image:
    rgb = np.asarray(image, dtype=np.uint8)
    lab = cv2.cvtColor(rgb, cv2.COLOR_RGB2LAB).astype(np.float32)
    current_mean, current_std = extraoral_lab_stats(image)
    adjusted = (lab - current_mean) * (target_std / current_std) + target_mean
    adjusted[:, :, 0] = np.clip(adjusted[:, :, 0], 0, 255)
    adjusted[:, :, 1] = np.clip(adjusted[:, :, 1], 0, 255)
    adjusted[:, :, 2] = np.clip(adjusted[:, :, 2], 0, 255)
    matched = cv2.cvtColor(adjusted.astype(np.uint8), cv2.COLOR_LAB2RGB)
    result = Image.fromarray(matched, "RGB")
    result = ImageEnhance.Contrast(result).enhance(1.02)
    result = ImageEnhance.Color(result).enhance(0.98)
    return ImageEnhance.Sharpness(result).enhance(1.03)


def extraoral_skin_mask(rgb: np.ndarray) -> np.ndarray:
    height, width, _ = rgb.shape
    r = rgb[:, :, 0]
    g = rgb[:, :, 1]
    b = rgb[:, :, 2]
    luminance = r * 0.299 + g * 0.587 + b * 0.114
    saturation = rgb.max(axis=2) - rgb.min(axis=2)
    mask = (
        (luminance > 80)
        & (luminance < 245)
        & (saturation > 5)
        & (saturation < 105)
        & (r > b + 5)
        & (g > b - 4)
        & (r > g * 0.86)
    )
    center = np.zeros((height, width), dtype=bool)
    center[int(height * 0.16):int(height * 0.78), int(width * 0.16):int(width * 0.84)] = True
    return mask & center


def rgb_stats(image: Image.Image, category: str = "") -> tuple[np.ndarray, np.ndarray]:
    arr = np.asarray(image, dtype=np.float32)
    pixels = arr.reshape(-1, 3)
    luminance = pixels[:, 0] * 0.299 + pixels[:, 1] * 0.587 + pixels[:, 2] * 0.114
    saturation = pixels.max(axis=1) - pixels.min(axis=1)
    if category == "extraoral":
        mask = (luminance > 70) & (luminance < 245) & (saturation < 125)
    else:
        low, high = np.percentile(luminance, [2, 98])
        mask = (luminance >= low) & (luminance <= high)
    sample = pixels[mask] if np.count_nonzero(mask) > 200 else pixels
    mean = sample.mean(axis=0)
    std = np.maximum(1.0, sample.std(axis=0))
    return mean, std


def match_rgb_stats(image: Image.Image, target_mean: np.ndarray, target_std: np.ndarray, category: str = "") -> Image.Image:
    arr = np.asarray(image, dtype=np.float32)
    current_mean, current_std = rgb_stats(image, category)
    adjusted = (arr - current_mean) * (target_std / current_std) + target_mean
    adjusted = np.clip(adjusted, 0, 255).astype(np.uint8)
    result = Image.fromarray(adjusted, "RGB")
    result = ImageEnhance.Contrast(result).enhance(1.03)
    result = ImageEnhance.Color(result).enhance(0.99)
    return ImageEnhance.Sharpness(result).enhance(1.04)


def fit_image(path: Path, target_w, target_h) -> Image.Image:
    target_ratio = target_w / target_h
    with Image.open(path) as image:
        image = ImageOps.exif_transpose(image).convert("RGB")
        ratio = image.width / image.height
        if ratio > target_ratio:
            new_w = int(image.height * target_ratio)
            left = (image.width - new_w) // 2
            image = image.crop((left, 0, left + new_w, image.height))
        else:
            new_h = int(image.width / target_ratio)
            top = (image.height - new_h) // 2
            image = image.crop((0, top, image.width, top + new_h))
        return image.resize((1600, 1060), Image.Resampling.LANCZOS)


def main() -> None:
    if not CERT_FILE.exists() or not KEY_FILE.exists():
        raise SystemExit(f"Missing HTTPS certificate: {CERT_FILE} / {KEY_FILE}")
    redirect_server = ThreadingHTTPServer(("0.0.0.0", HTTP_PORT), RedirectHandler)
    threading.Thread(target=redirect_server.serve_forever, daemon=True).start()

    server = ThreadingHTTPServer(("0.0.0.0", HTTPS_PORT), DentalHandler)
    context = ssl.SSLContext(ssl.PROTOCOL_TLS_SERVER)
    context.load_cert_chain(certfile=str(CERT_FILE), keyfile=str(KEY_FILE))
    server.socket = context.wrap_socket(server.socket, server_side=True)
    print(f"Dental PPT listening on https://0.0.0.0:{HTTPS_PORT}", flush=True)
    server.serve_forever()


if __name__ == "__main__":
    main()
